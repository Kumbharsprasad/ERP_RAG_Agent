import io
import os
from bs4 import BeautifulSoup, NavigableString
import pandas as pd
import logfire

def parse_with_liteparse(file_bytes: bytes, filename: str, ext: str) -> tuple[list[dict], list[str]]:
    """
    Unified parser using LiteParse for PDF, DOCX, PPTX, XLSX, and images (PNG/JPG/JPEG).
    Runs complexity check for PDF and images first to dynamically choose OCR setting.
    """
    from liteparse import LiteParse
    
    # Decide if we run complexity check (only for PDF and images)
    is_pdf_or_image = ext.lower() in [".pdf", ".png", ".jpg", ".jpeg"]
    ocr_needed = True
    
    if is_pdf_or_image:
        try:
            temp_parser = LiteParse()
            pages_complexity = temp_parser.is_complex(file_bytes)
            ocr_needed = False
            for p_comp in pages_complexity:
                needs_ocr = getattr(p_comp, "needs_ocr", False)
                page_num = getattr(p_comp, "page_num", 0)
                reasons = getattr(p_comp, "reasons", [])
                logfire.info(
                    "Page {page_num} complexity check: needs_ocr={needs_ocr}, reasons={reasons}",
                    page_num=page_num, needs_ocr=needs_ocr, reasons=reasons
                )
                if needs_ocr:
                    ocr_needed = True
        except Exception as e:
            logfire.warn("Complexity check failed: {error}. Defaulting to ocr_enabled=True", error=str(e))
            ocr_needed = True
    else:
        # Non-PDFs/images (like DOCX, PPTX, XLSX) do not require OCR for standard text layers
        ocr_needed = False
        
    logfire.info(
        "Parsing {filename} via LiteParse with ocr_enabled={ocr_enabled}",
        filename=filename, ocr_enabled=ocr_needed
    )
    
    try:
        parser = LiteParse(output_format="markdown", ocr_enabled=ocr_needed, image_mode="placeholder")
        result = parser.parse(file_bytes)
    except Exception as e:
        err_msg = str(e)
        is_missing_dependency = (
            "LibreOffice is not installed" in err_msg or 
            "ImageMagick is not installed" in err_msg or 
            "conversion error" in err_msg
        )
        if is_missing_dependency:
            missing_dep = "LibreOffice"
            if "ImageMagick" in err_msg:
                missing_dep = "ImageMagick"
            logfire.warn(
                "LiteParse parsing failed due to missing system dependency: {missing_dependency}. Error: {error}. Attempting local library fallback...",
                missing_dependency=missing_dep,
                error=err_msg
            )
            if ext.lower() == ".docx":
                try:
                    import docx
                    doc = docx.Document(io.BytesIO(file_bytes))
                    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
                    blocks = []
                    chunk_size = 10
                    for i in range(0, len(paragraphs), chunk_size):
                        group = paragraphs[i:i + chunk_size]
                        blocks.append({
                            "text": "\n".join(group),
                            "source_file": filename,
                            "location": f"Paragraphs {i + 1}-{i + len(group)}"
                        })
                    logfire.info("Successfully parsed {filename} via fallback python-docx.", filename=filename)
                    warning_msg = "PPTX/DOCX parsing used degraded fallback (python-pptx/python-docx) instead of LiteParse due to missing system dependency: LibreOffice. Parsing quality may be reduced for tables and complex layouts."
                    return blocks, [warning_msg]
                except Exception as docx_err:
                    logfire.error("Fallback docx parsing failed: {error}", error=str(docx_err))
            elif ext.lower() == ".pptx":
                try:
                    from pptx import Presentation
                    prs = Presentation(io.BytesIO(file_bytes))
                    blocks = []
                    for idx, slide in enumerate(prs.slides):
                        slide_text_parts = []
                        for shape in slide.shapes:
                            if hasattr(shape, "text") and shape.text.strip():
                                slide_text_parts.append(shape.text)
                        blocks.append({
                            "text": "\n".join(slide_text_parts),
                            "source_file": filename,
                            "location": f"Slide {idx + 1}"
                        })
                    logfire.info("Successfully parsed {filename} via fallback python-pptx.", filename=filename)
                    warning_msg = "PPTX/DOCX parsing used degraded fallback (python-pptx/python-docx) instead of LiteParse due to missing system dependency: LibreOffice. Parsing quality may be reduced for tables and complex layouts."
                    return blocks, [warning_msg]
                except Exception as pptx_err:
                    logfire.error("Fallback pptx parsing failed: {error}", error=str(pptx_err))
            elif ext.lower() == ".xlsx":
                try:
                    xls = pd.ExcelFile(io.BytesIO(file_bytes))
                    blocks = []
                    for sheet_name in xls.sheet_names:
                        df = pd.read_excel(xls, sheet_name=sheet_name)
                        blocks.append({
                            "text": df.to_csv(index=False).strip(),
                            "source_file": filename,
                            "location": f"Sheet: {sheet_name}"
                        })
                    logfire.info("Successfully parsed {filename} via fallback pandas.read_excel.", filename=filename)
                    return blocks, []
                except Exception as xlsx_err:
                    logfire.error("Fallback xlsx parsing failed: {error}", error=str(xlsx_err))
        raise e
        
    blocks = []
    
    # Safely check if result contains pages and they have the 'text' property populated
    has_page_text = False
    if hasattr(result, "pages") and result.pages:
        first_page = result.pages[0]
        if hasattr(first_page, "text") and getattr(first_page, "text") is not None:
            has_page_text = True
            
    if has_page_text:
        for idx, page in enumerate(result.pages):
            p_num = getattr(page, "page_num", idx + 1)
            text_content = getattr(page, "text", "")
            
            if ext.lower() == ".pptx":
                location = f"Slide {p_num}"
            elif ext.lower() == ".xlsx":
                sheet_name = getattr(page, "sheet_name", None)
                if sheet_name:
                    location = f"Sheet: {sheet_name}"
                else:
                    location = f"Sheet {p_num}"
            else:
                location = f"Page {p_num}"
                
            blocks.append({
                "text": text_content,
                "source_file": filename,
                "location": location
            })
    else:
        # Fallback: split the whole document text by markdown page breaks
        full_text = getattr(result, "text", "") or ""
        page_texts = full_text.split('\f')
        if len(page_texts) <= 1:
            page_texts = full_text.split('\n---\n')
            
        for idx, text_content in enumerate(page_texts):
            p_num = idx + 1
            if ext.lower() == ".pptx":
                location = f"Slide {p_num}"
            elif ext.lower() == ".xlsx":
                location = f"Sheet {p_num}"
            else:
                location = f"Page {p_num}"
                
            blocks.append({
                "text": text_content.strip(),
                "source_file": filename,
                "location": location
            })
            
    return blocks, []


def parse_html(file_bytes: bytes, filename: str) -> list[dict]:
    """
    Parses an HTML file from bytes.
    Extracts visible text, splitting content into sections based on heading tags (h1-h6).
    Location is formatted as 'Section {n}' (1-indexed).
    """
    html_content = file_bytes.decode("utf-8", errors="ignore")
    soup = BeautifulSoup(html_content, "html.parser")
    
    # Remove scripts, styles, and other metadata tags to clean the HTML
    for tag in soup(["script", "style", "meta", "noscript", "header", "footer"]):
        tag.decompose()
        
    blocks = []
    current_section_text = []
    section_idx = 1
    
    def traverse(node):
        nonlocal section_idx, current_section_text
        # If node is one of the heading tags
        if node.name in ["h1", "h2", "h3", "h4", "h5", "h6"]:
            # Save the text gathered so far into a block
            text_so_far = "".join(current_section_text).strip()
            if text_so_far:
                blocks.append({
                    "text": text_so_far,
                    "source_file": filename,
                    "location": f"Section {section_idx}"
                })
                section_idx += 1
                current_section_text = []
            
            # Start the new section with the heading text itself
            heading_text = node.get_text().strip()
            if heading_text:
                current_section_text.append(heading_text + "\n")
            return
            
        # If node is a text node (NavigableString), collect it
        if isinstance(node, NavigableString):
            val = str(node).strip()
            if val:
                current_section_text.append(val + " ")
            return
            
        # Recursively traverse child nodes
        if hasattr(node, "children"):
            for child in node.children:
                traverse(child)
 
    # Begin traversal on body (or root soup if body is missing)
    start_node = soup.body if soup.body else soup
    traverse(start_node)
    
    # Grab any trailing text
    final_text = "".join(current_section_text).strip()
    if final_text:
        blocks.append({
            "text": final_text,
            "source_file": filename,
            "location": f"Section {section_idx}"
        })
        
    return blocks
 
def parse_csv(file_bytes: bytes, filename: str) -> list[dict]:
    """
    Parses a CSV file from bytes.
    Groups rows into blocks of 20 rows each, including headers in every block.
    Location is formatted as 'Rows {start}-{end}'.
    """
    df = pd.read_csv(io.BytesIO(file_bytes))
    blocks = []
    chunk_size = 20
    num_rows = len(df)
    
    if num_rows == 0:
        blocks.append({
            "text": ",".join(df.columns),
            "source_file": filename,
            "location": "Rows 0-0"
        })
        return blocks
        
    for i in range(0, num_rows, chunk_size):
        sub_df = df.iloc[i:i + chunk_size]
        text = sub_df.to_csv(index=False)
        start_row = i + 1
        end_row = min(i + len(sub_df), num_rows)
        blocks.append({
            "text": text.strip(),
            "source_file": filename,
            "location": f"Rows {start_row}-{end_row}"
        })
        
    return blocks
 
def parse_txt(file_bytes: bytes, filename: str) -> list[dict]:
    """
    Parses a TXT file from bytes.
    Splits the decoded text into blocks of ~500 words.
    Location is formatted as 'Block {n}'.
    """
    text = file_bytes.decode("utf-8", errors="ignore")
    words = text.split()
    blocks = []
    chunk_size = 500
    num_words = len(words)
    
    if num_words == 0:
        blocks.append({
            "text": "",
            "source_file": filename,
            "location": "Block 1"
        })
        return blocks
        
    for i in range(0, num_words, chunk_size):
        group = words[i:i + chunk_size]
        blocks.append({
            "text": " ".join(group),
            "source_file": filename,
            "location": f"Block {i // chunk_size + 1}"
        })
        
    return blocks
 
def parse_file(file_bytes: bytes, filename: str) -> tuple[list[dict], list[str]]:
    """
    Dispatcher function that picks the right parser based on file extension.
    Raises ValueError for unsupported types.
    """
    ext = os.path.splitext(filename)[1].lower()
    is_liteparse = ext in [".pdf", ".docx", ".pptx", ".xlsx", ".png", ".jpg", ".jpeg"]
    parser_name = "LiteParse" if is_liteparse else "Custom"
    
    with logfire.span("Parsing file {filename} ({ext}) via {parser_name}", filename=filename, ext=ext, parser_name=parser_name):
        if is_liteparse:
            blocks, warnings = parse_with_liteparse(file_bytes, filename, ext)
        elif ext in [".html", ".htm"]:
            blocks = parse_html(file_bytes, filename)
            warnings = []
        elif ext == ".csv":
            blocks = parse_csv(file_bytes, filename)
            warnings = []
        elif ext == ".txt":
            blocks = parse_txt(file_bytes, filename)
            warnings = []
        else:
            logfire.error("Unsupported extension {ext} for {filename}", ext=ext, filename=filename)
            raise ValueError(f"Unsupported file type: '{ext}' in file '{filename}'")
            
        logfire.info("Successfully parsed {filename} via {parser_name}. Extracted {num_blocks} blocks.", filename=filename, parser_name=parser_name, num_blocks=len(blocks))
        return blocks, warnings
